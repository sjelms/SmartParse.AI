# smartparse_watch.py
# This script monitors a folder and its subfolders for new files (PDFs or images)
# and routes them to appropriate processing functions.

import base64
import os
import sys
import threading
import time
from datetime import datetime
from pathlib import Path
from typing import Dict, Optional, Tuple

import queue

from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

from dotenv import load_dotenv
import openai
from bs4 import BeautifulSoup
import fitz  # PyMuPDF

try:
    from pync import Notifier
    notifier_available = True
except ImportError:
    notifier_available = False
    Notifier = None  # type: ignore

from macos_tags import add as add_finder_tag, Tag, Color, set_all  # type: ignore

# Load environment variables from .env file
load_dotenv()


def _resolve_api_key() -> str:
    """Return a usable OpenAI API key or exit with an actionable message."""

    candidate: Optional[str] = None
    try:
        candidate = openai.api_key or None  # type: ignore[attr-defined]
    except Exception:
        candidate = None

    if not candidate:
        candidate = os.getenv("OPENAI_API_KEY")

    if candidate:
        candidate = candidate.strip()

    if not candidate:
        raise SystemExit(
            "OPENAI_API_KEY is not set. Export it before starting SmartParse or run the command through `op run --env-file=.env -- ...`."
        )

    if candidate.startswith("op://"):
        raise SystemExit(
            "OPENAI_API_KEY resolves to a 1Password secret reference (op://...). Unlock 1Password (`op signin`) and launch SmartParse via `op run --env-file=.env -- ...` so the real key is available."
        )

    if candidate.lower().startswith("your-") or "api-key" in candidate.lower():
        raise SystemExit(
            "OPENAI_API_KEY is still set to a placeholder value. Replace it with your real OpenAI key before running SmartParse."
        )

    return candidate


api_key = _resolve_api_key()
openai.api_key = api_key
client = openai.OpenAI(api_key=api_key)

# Define which model to use for each type
MODEL_TEXT = "gpt-3.5-turbo"
MODEL_PDF = "gpt-3.5-turbo"
MODEL_IMAGE = "gpt-4o-2024-11-20"

def _read_simple_config(path: Path) -> Dict[str, str]:
    data: Dict[str, str] = {}
    for raw_line in path.read_text(encoding="utf-8").splitlines():
        line = raw_line.split("#", 1)[0].strip()
        if not line or ":" not in line:
            continue
        key, value = line.split(":", 1)
        cleaned = value.strip()
        if cleaned and cleaned[0] == cleaned[-1] and cleaned[0] in ('"', '\''):
            cleaned = cleaned[1:-1]
        data[key.strip()] = cleaned
    return data


def _expand_path(raw_value: str, base: Path) -> Path:
    expanded = os.path.expandvars(raw_value)
    expanded = os.path.expanduser(expanded)
    candidate = Path(expanded)
    if not candidate.is_absolute():
        candidate = (base / candidate).resolve()
    return candidate


def _load_config() -> Tuple[Dict[str, str], Optional[Path]]:
    search_candidates = []
    env_path = os.getenv("SMARTPARSE_CONFIG")
    if env_path:
        search_candidates.append(Path(env_path).expanduser())
    repo_config = Path(__file__).resolve().parent / "config.yaml"
    search_candidates.append(repo_config)
    home_config = Path.home() / ".config" / "smartparse" / "config.yaml"
    search_candidates.append(home_config)

    for candidate in search_candidates:
        if candidate.is_file():
            try:
                return _read_simple_config(candidate), candidate
            except Exception as exc:
                print(f"Failed to read config {candidate}: {exc}", file=sys.stderr)
    return {}, None


CONFIG, CONFIG_PATH = _load_config()
CONFIG_BASE = CONFIG_PATH.parent if CONFIG_PATH else Path(__file__).resolve().parent


if len(sys.argv) > 1:
    WATCH_DIR = Path(sys.argv[1]).expanduser().resolve()
else:
    watch_from_config = CONFIG.get("watch_directory")
    if watch_from_config:
        WATCH_DIR = _expand_path(watch_from_config, CONFIG_BASE)
    else:
        WATCH_DIR = Path.home() / "Library" / "Mobile Documents" / "com~apple~CloudDocs" / "Documents" / "SmartParseWatch"

# File extensions grouped by type
IMAGE_EXTENSIONS = {
    ".jpg", ".jpeg", ".png", ".webp",
    ".heic", ".heif", ".gif", ".bmp", ".tiff", ".tif"
}
TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".html", ".htm", ".xhtml"
}
PDF_EXTENSION = ".pdf"
OFFICE_EXTENSIONS = {
    ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx"
}

# Max files to queue for processing at a time
MAX_QUEUE_SIZE = 20
file_queue: queue.Queue = queue.Queue(maxsize=MAX_QUEUE_SIZE)

import json

logs_from_config = CONFIG.get("logs_dir")
if logs_from_config:
    LOG_DIR = _expand_path(logs_from_config, CONFIG_BASE)
else:
    LOG_DIR = Path(__file__).resolve().parent / "logs"

RUN_ID = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
LOG_FILE = LOG_DIR / f"smartparse_log_{RUN_ID}.jsonl"

# Ensure log directory exists
def ensure_log_dir():
    LOG_DIR.mkdir(parents=True, exist_ok=True)

# Log operation
def log_file_operation(original_path, new_path, file_type, tag, status, error=None):
    ensure_log_dir()
    entry = {
        "timestamp": time.strftime("%Y-%m-%d %H:%M:%S", time.localtime()),
        "original_filename": str(original_path),
        "new_filename": str(new_path) if new_path else None,
        "file_type": file_type,
        "tag": tag,
        "status": status,
        "error": error,
    }
    with open(LOG_FILE, "a", encoding="utf-8") as f:
        f.write(json.dumps(entry) + "\n")

def notify_user(message: str, title: str = "SmartParse.AI", sound: str = "default") -> None:
    """
    Sends a user notification if possible. Falls back to print if no GUI notifier is available.
    """
    if notifier_available and Notifier is not None:
        try:
            Notifier.notify(message, title=title, sound=sound)
        except Exception:
            print(f"Notification failed: {message}")
    else:
        print(f"{title}: {message}")

def worker() -> None:
    """
    Background worker thread that processes files from the queue.
    """
    while True:
        filepath = file_queue.get()
        if filepath is None:
            break
        try:
            # --- Race Condition Handling ---
            # Wait briefly and check if file size remains stable
            time.sleep(1)
            initial_size = filepath.stat().st_size
            time.sleep(1)
            final_size = filepath.stat().st_size

            if initial_size != final_size:
                print(f"File size still changing for {filepath.name}. Re-queuing.")
                file_queue.put(filepath)  # Re-queue file for later
                file_queue.task_done()
                continue  # Skip to next file

            FileHandler().handle_file(filepath)
            file_queue.task_done()
            if file_queue.empty():
                refill_queue()
        except Exception as e:
            print(f"Error processing file {filepath}: {e}")
            file_queue.task_done()


# Start the background worker thread
threading.Thread(target=worker, daemon=True).start()


def refill_queue() -> None:
    """
    Scans the WATCH_DIR for any remaining unprocessed files and queues the next batch.
    Skips files already renamed or moved.
    """
    print("Scanning for additional files to queue...")
    try:
        files = list(WATCH_DIR.glob("*"))
        queued = 0
        for filepath in files:
            if (
                not filepath.is_file()
                or filepath.name.startswith("failed_")
                or filepath.name.startswith(".")  # Ignore hidden files like .DS_Store
                or filepath.parent != WATCH_DIR
                or queued >= MAX_QUEUE_SIZE
            ):
                continue
            try:
                file_queue.put_nowait(filepath)
                print(f"Re-queued file: {filepath}")
                queued += 1
            except queue.Full:
                print("Queue is full. Refill paused.")
                break
    except Exception as e:
        print(f"Error during refill queue scan: {e}")

def mark_as_failed(filepath: Path) -> Path:
    """
    Marks a file as failed by renaming it with a 'failed_' prefix in the same directory.
    This function is used to flag files that could not be processed successfully.
    """
    parent = filepath.parent
    failed_name = f"failed_{filepath.name}"
    failed_path = parent / failed_name
    try:
        filepath.rename(failed_path)
        print(f"Marked file as failed: {failed_path}")
        notify_user(f"Could not process: {filepath.name}")
    except Exception as e:
        print(f"Could not mark file as failed ({filepath}): {e}")
    return failed_path

def move_file_to_subfolder(filepath: Path, folder_name: str) -> Path:
    """
    Moves the given file to a specified subfolder within the WATCH_DIR.
    Creates the subfolder if it does not exist.
    """
    target_folder = WATCH_DIR / folder_name
    target_folder.mkdir(parents=True, exist_ok=True)
    target_path = target_folder / filepath.name
    filepath.rename(target_path)
    print(f"Moved file to: {target_path}")
    return target_path

# Handler class that responds to file system events such as file creation.
class FileHandler(FileSystemEventHandler):
    def on_created(self, event) -> None:
        """
        Called when a file or directory is created.
        """
        if event.is_directory:
            return
        src_path = event.src_path.decode() if isinstance(event.src_path, bytes) else event.src_path
        filepath = Path(src_path)
        print(f"Queued file: {filepath}")
        try:
            file_queue.put_nowait(filepath)
        except queue.Full:
            print(f"Queue full. Skipping file: {filepath}")

    def process_image(self, filepath: Path) -> None:
        """
        Processes image files using GPT-4o vision model to generate a descriptive filename and Finder tag.
        Renames and moves the file to the 'images' subfolder.
        """
        print(f"Processing image: {filepath}")
        try:
            with open(filepath, "rb") as img_file:
                image_data = img_file.read()
                image_base64 = base64.b64encode(image_data).decode("utf-8")

            allowed_categories = [
                "Quote", "Sign", "Cartoon", "Meme", "Photo", "Illustration", "Diagram", "Screenshot", "Logo", "Map", "Chart/Graph"
            ]
            messages = [
                {
                    "role": "system",
                    "content": (
                        "You are a helpful assistant that generates descriptive, searchable filenames and a category label based on the visual and textual content of an image. "
                        "Detect whether the image contains visible text. If text is present, extract key phrases and any clear sentiment or tone (e.g., urgent, humorous, inspirational) and incorporate them succinctly. "
                        "By default, the filename must consist of 10 to 12 lowercase words, no punctuation, no underscores, no file extension. If the image contains text, you may use up to 16 lowercase words to capture key text and sentiment. Be concise but informative. "
                        "Return ONLY a JSON object with two fields: 'description' (10–12 lowercase words, or up to 16 if the image contains text; no punctuation, no underscores, no file extension) and 'category' (one of: Quote, Sign, Cartoon, Meme, Photo, Illustration, Diagram, Screenshot, Logo, Map, Chart/Graph). "
                        "The 'description' will be used as the filename (with a timestamp), and the 'category' will be used as a Finder tag (not in the filename). "
                        "If the content does not fit any category, use 'Other'. "
                        "Example without text: {\"description\": \"golden retriever surfing ocean wave at sunset\", \"category\": \"Photo\"} "
                        "Example with text: {\"description\": \"conference badge speaker anna lee ai ethics keynote optimistic\", \"category\": \"Screenshot\"}"
                    ),
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/{filepath.suffix[1:]};base64,{image_base64}"
                            },
                        }
                    ],
                },
            ]
            response = client.chat.completions.create(
                model=MODEL_IMAGE,
                messages=messages,
                max_tokens=80,
                temperature=0.2,
                n=1,
            )
            try:
                result = json.loads((response.choices[0].message.content or '').strip())
                description = (result.get('description') or '').lower().replace('_', ' ').strip()
                category = (result.get('category') or 'Other')
                if not isinstance(category, str):
                    category = str(category)
                category = (category or '').strip()
                # Validate description: non-empty, no braces/slashes, and not excessively long
                if not description or any(c in description for c in '{}[]/\\') or len(description) > 160:
                    print(f"Invalid description from AI: {description}. Marking as failed.")
                    log_file_operation(filepath, None, "image", None, "fail", error="Invalid description from AI")
                    mark_as_failed(filepath)
                    return
            except Exception:
                print(f"Failed to parse AI response as JSON for {filepath}. Marking as failed.")
                log_file_operation(filepath, None, "image", None, "fail", error="Failed to parse AI response as JSON")
                mark_as_failed(filepath)
                return
            timestamp = get_file_datetime_string(filepath)
            new_name = f"{description}_{timestamp}{filepath.suffix.lower()}"
            new_path = filepath.with_name(new_name)
            filepath.rename(new_path)
            print(f"Renamed image to: {new_path}")
            # Use string method to set tag color for new tags
            tag = f"{category}\n{Color.YELLOW}"
            set_all([tag], file=str(new_path))  # Use set_all to ensure color is applied
            moved_path = move_file_to_subfolder(new_path, "images")
            notify_user("Image processed and renamed")
            log_file_operation(filepath, new_path, "image", category, "success")
        except Exception as e:
            print(f"Error processing image file {filepath}: {e}")
            log_file_operation(filepath, None, "image", None, "fail", error=str(e))
            mark_as_failed(filepath)

    def process_pdf(self, filepath: Path) -> None:
        """
        Processes PDF files by extracting text to generate a descriptive filename and Finder tag.
        Renames and moves the file to the 'pdfs' subfolder.
        """
        print(f"Processing PDF: {filepath}")
        try:
            doc = fitz.open(filepath)
            if doc.page_count > 0:
                try:
                    page = doc.load_page(0)
                    text = page.get_text().strip()  # type: ignore
                except Exception as e:
                    print(f"Error extracting text from first page: {e}")
                    text = ""
            else:
                text = ""
            if not text:
                print(f"PDF {filepath} has no extractable text. Marking as failed.")
                log_file_operation(filepath, None, "pdf", None, "fail", error="No extractable text")
                mark_as_failed(filepath)
                return
            allowed_categories = [
                "Book", "Academic Paper", "Contract", "Invoice", "Receipt", "Statement", "Manual/Guide", "Report", "Form", "Presentation", "Brochure/Flyer", "Resume/CV", "Letter", "Certificate", "Agreement"
            ]
            description, category = generate_filename_and_category_from_text(
                text,
                model=MODEL_PDF,
                allowed_categories=allowed_categories,
                prompt_extra="Prioritize including key identifiers such as paper titles, author names, organizations (e.g. McKinsey, Ministry of Housing, Columbia University), or publication dates. The filename must consist of between 7 and 10 lowercase words, no punctuation, no underscores, no file extension. Be as descriptive as possible within this range."
            )
            timestamp = get_file_datetime_string(filepath)
            new_name = f"{description}_{timestamp}{filepath.suffix.lower()}"
            new_path = filepath.with_name(new_name)
            filepath.rename(new_path)
            print(f"Renamed PDF to: {new_path}")
            # Use string method to set tag color for new tags
            tag = f"{category}\n{Color.RED}"
            set_all([tag], file=str(new_path))  # Use set_all to ensure color is applied
            moved_path = move_file_to_subfolder(new_path, "pdfs")
            notify_user("PDF processed and renamed")
            log_file_operation(filepath, new_path, "pdf", category, "success")
        except Exception as e:
            print(f"Error processing PDF file {filepath}: {e}")
            log_file_operation(filepath, None, "pdf", None, "fail", error=str(e))
            mark_as_failed(filepath)

    def process_textfile(self, filepath: Path) -> None:
        """
        Processes text files by reading their content, generating a descriptive filename and Finder tag,
        renaming, and moving them to the 'text' subfolder.
        Supports plain text, HTML, and Microsoft Office files.
        """
        ext = filepath.suffix.lower()
        try:
            text = ""
            if ext in {".html", ".htm", ".xhtml"}:
                with open(filepath, "r", encoding="utf-8") as f:
                    soup = BeautifulSoup(f, "html.parser")
                    text = soup.get_text(separator=' ', strip=True)
            elif ext == ".docx":
                from docx import Document  # type: ignore
                doc = Document(filepath)
                text = '\n'.join([para.text for para in doc.paragraphs])
            elif ext == ".xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(filepath, data_only=True)
                lines = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        lines.append(' '.join([str(cell) if cell is not None else '' for cell in row]))
                text = '\n'.join(lines)
            elif ext == ".pptx":
                from pptx import Presentation  # type: ignore
                prs = Presentation(filepath)
                lines = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            lines.append(shape.text)
                text = '\n'.join(lines)
            elif ext in {".doc", ".xls", ".ppt"}:
                print(f"Legacy Office format not supported for text extraction: {filepath}")
                mark_as_failed(filepath)
                return
            else:
                with open(filepath, "r", encoding="utf-8") as f:
                    text = f.read()
            if not text.strip():
                print(f"Text file {filepath} has insufficient content. Marking as failed.")
                log_file_operation(filepath, None, "text", None, "fail", error="Insufficient content")
                mark_as_failed(filepath)
                return
            allowed_categories = [
                "Notes", "Outline", "Draft", "Paper", "Journal Entry", "List", "Code", "Markdown", "Recipe", "Correspondence", "Brainstorm", "Transcript"
            ]
            description, category = generate_filename_and_category_from_text(
                text,
                model=MODEL_TEXT,
                allowed_categories=allowed_categories,
                prompt_extra=(
                    "The filename must consist of between 7 and 10 lowercase words, no punctuation, no underscores, no file extension. Be as descriptive as possible within this range. "
                    "Choose the most specific and accurate category from the list based on the content. "
                    "Do not default to 'Correspondence' unless the file is truly a letter, email, or message. "
                    "If the file is a transcript of a conversation, interview, or meeting, use 'Transcript'. "
                    "If the content is a list, use 'List'. If it is a draft, use 'Draft', etc."
                )
            )
            timestamp = get_file_datetime_string(filepath)
            new_name = f"{description}_{timestamp}{ext}"
            new_path = filepath.with_name(new_name)
            filepath.rename(new_path)
            print(f"Renamed file to: {new_path}")
            # Use string method to set tag color for new tags
            tag = f"{category}\n{Color.BLUE}"
            set_all([tag], file=str(new_path))  # Use set_all to ensure color is applied
            moved_path = move_file_to_subfolder(new_path, "text")
            notify_user("Text file processed and renamed")
            log_file_operation(filepath, new_path, "text", category, "success")
        except Exception as e:
            print(f"Error processing text file {filepath}: {e}")
            log_file_operation(filepath, None, "text", None, "fail", error=str(e))
            mark_as_failed(filepath)

    def handle_file(self, filepath: Path) -> None:
        """
        Determines the file type based on extension and routes to the correct processor.
        """
        ext = filepath.suffix.lower()
        if ext == PDF_EXTENSION:
            self.process_pdf(filepath)
        elif ext in IMAGE_EXTENSIONS:
            self.process_image(filepath)
        elif ext in TEXT_EXTENSIONS or ext in OFFICE_EXTENSIONS:
            self.process_textfile(filepath)
        else:
            print(f"Unsupported file type: {filepath.name}")
            log_file_operation(filepath, None, "unknown", None, "fail", error="Unsupported file type")
            mark_as_failed(filepath)

def generate_filename_and_category_from_text(text: str, model: str, allowed_categories: list[str], prompt_extra: str = "") -> tuple[str, str]:
    """
    Uses OpenAI's API to generate a short, descriptive filename and a category based on provided text.
    Returns (description, category). Only categories from allowed_categories are valid.
    """
    allowed_str = ', '.join(allowed_categories)
    system_prompt = (
        f"You are a helpful assistant that generates short, descriptive filenames and a category label based on provided file content. "
        f"Return ONLY a JSON object with two fields: 'description' (between 7 and 10 lowercase words, no punctuation, no underscores, no file extension, be as descriptive as possible within this range) and 'category' (one of: {allowed_str}). "
        f"The 'description' will be used as the filename (with a timestamp), and the 'category' will be used as a Finder tag (not in the filename). "
        f"If the content does not fit any category, use 'Other'. "
        f"Example: {{\"description\": \"product configuration in construction patrik jensen\", \"category\": \"Report\"}} (filename: product configuration in construction patrik jensen_2025-07-19_13.05.10.pdf, tag: Report)"
    )
    if prompt_extra:
        system_prompt += " " + prompt_extra
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": text},
    ]
    response = client.chat.completions.create(
        model=model,
        messages=messages,  # type: ignore
        max_tokens=80,
        temperature=0.2,
        n=1,
    )
    import json
    try:
        result = json.loads((response.choices[0].message.content or '').strip())
        description = (result.get('description') or '').lower().replace('_', ' ').strip()
        category = (result.get('category') or 'Other').strip()
    except Exception:
        # fallback: use all as description, category Other
        description = (response.choices[0].message.content or '').strip().lower().replace('_', ' ').strip()
        category = 'Other'
    return description, category

def get_file_datetime_string(file_path: Path) -> str:
    """
    Returns the file's creation time (preferred) or modification time if creation time is unavailable.
    The timestamp is formatted as 'YYYY-MM-DD_HH.MM.SS'.
    """
    try:
        # Try to get the creation time (macOS)
        timestamp = file_path.stat().st_birthtime
    except AttributeError:
        # Fallback to modified time (cross-platform)
        timestamp = file_path.stat().st_mtime

    # Format the timestamp to 'YYYY-MM-DD_HH.MM.SS'
    return time.strftime("%Y-%m-%d_%H.%M.%S", time.localtime(timestamp))

def show_final_dialog(success_count, issue_count, fail_count, image_count, pdf_count, text_count, elapsed_sec):
    import subprocess
    time_min = int(elapsed_sec // 60)
    time_sec = int(elapsed_sec % 60)
    time_str = f"{time_min}m {time_sec}s" if time_min else f"{time_sec}s"
    applescript = f'''
    display dialog "SmartParse.AI Batch Summary

    ✅ Files processed successfully: {success_count}
    ⚠️ Files with issues: {issue_count}
    ❌ Files failed: {fail_count}

    🖼️ Images: {image_count}
    📄 PDFs: {pdf_count}
    📝 Text: {text_count}

    🕒 Time elapsed: {time_str}
    " buttons {{"OK"}} default button "OK" with title "SmartParse.AI"
    '''
    subprocess.run(['osascript', '-e', applescript])

# Main execution block to start the folder watcher
if __name__ == "__main__":
    print(f"Running SmartParse.AI batch processor on: {WATCH_DIR}")

    import time as _time
    start_time = _time.time()

    # In-memory stats for this batch only
    batch_stats = {
        'image_count': 0,
        'pdf_count': 0,
        'text_count': 0,
        'fail_count': 0,
    }

    # Patch FileHandler to update batch_stats
    orig_process_image = FileHandler.process_image
    orig_process_pdf = FileHandler.process_pdf
    orig_process_textfile = FileHandler.process_textfile
    orig_mark_as_failed = mark_as_failed

    def patched_process_image(self, filepath):
        try:
            orig_process_image(self, filepath)
            if not filepath.name.startswith("failed_"):
                batch_stats['image_count'] += 1
        except Exception:
            batch_stats['fail_count'] += 1
            raise
    def patched_process_pdf(self, filepath):
        try:
            orig_process_pdf(self, filepath)
            if not filepath.name.startswith("failed_"):
                batch_stats['pdf_count'] += 1
        except Exception:
            batch_stats['fail_count'] += 1
            raise
    def patched_process_textfile(self, filepath):
        try:
            orig_process_textfile(self, filepath)
            if not filepath.name.startswith("failed_"):
                batch_stats['text_count'] += 1
        except Exception:
            batch_stats['fail_count'] += 1
            raise
    def patched_mark_as_failed(filepath):
        batch_stats['fail_count'] += 1
        return orig_mark_as_failed(filepath)

    FileHandler.process_image = patched_process_image
    FileHandler.process_pdf = patched_process_pdf
    FileHandler.process_textfile = patched_process_textfile
    mark_as_failed = patched_mark_as_failed

    try:
        files_remaining = True
        while files_remaining:
            files_remaining = False
            files = list(WATCH_DIR.glob("*"))
            for filepath in files:
                if (
                    filepath.is_file()
                    and not filepath.name.startswith("failed_")
                    and filepath.parent == WATCH_DIR
                    and not filepath.name.startswith(".")
                ):
                    try:
                        file_queue.put_nowait(filepath)
                        print(f"Queued file: {filepath}")
                        files_remaining = True
                    except queue.Full:
                        print(f"Queue full. Pausing before next refill...")
                        break
            file_queue.join()

        # Only show stats for this batch
        success_count = batch_stats['image_count'] + batch_stats['pdf_count'] + batch_stats['text_count']
        issue_count = 0  # Can be expanded if you want to track partials
        fail_count = batch_stats['fail_count']
        image_count = batch_stats['image_count']
        pdf_count = batch_stats['pdf_count']
        text_count = batch_stats['text_count']
        elapsed_sec = int(_time.time() - start_time)

        show_final_dialog(success_count, issue_count, fail_count, image_count, pdf_count, text_count, elapsed_sec)

        file_queue.put(None)  # type: ignore  # Signal the worker thread to exit

    except Exception as e:
        print(f"Batch processing failed: {e}")
